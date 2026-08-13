from __future__ import annotations

import argparse
import hashlib
import json
import re
import zipfile
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


CJK_RE = re.compile(r"[\u3400-\u9fff]")


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def iter_shapes(shapes):
    """Yield top-level and grouped shapes without losing nested pictures."""
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--expected-font", default="Microsoft YaHei")
    parser.add_argument("--json-out", type=Path, required=True)
    args = parser.parse_args()

    with zipfile.ZipFile(args.pptx) as archive:
        corrupt_member = archive.testzip()
        names = archive.namelist()
        slide_xml = [name for name in names if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)]
        media = [name for name in names if name.startswith("ppt/media/") and not name.endswith("/")]
        raw_xml = "\n".join(
            archive.read(name).decode("utf-8", errors="replace")
            for name in slide_xml
        )

    deck = Presentation(str(args.pptx))
    text_shapes = 0
    picture_shapes = 0
    texts: list[str] = []
    fonts: Counter[str] = Counter()
    for slide in deck.slides:
        for shape in iter_shapes(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_shapes += 1
            if not getattr(shape, "has_text_frame", False):
                continue
            text_shapes += 1
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    texts.append(run.text)
                    if run.font.name:
                        fonts[run.font.name] += 1

    width = deck.slide_width
    height = deck.slide_height
    aspect = width / height
    cjk_count = sum(len(CJK_RE.findall(text)) for text in texts)
    checks = {
        "zip_crc_clean": corrupt_member is None,
        "python_pptx_open": True,
        "slide_count_at_least_3": len(deck.slides) >= 3,
        "aspect_ratio_16_9": abs(aspect - (16 / 9)) < 0.001,
        "editable_text_shapes_present": text_shapes >= 3,
        "editable_cjk_present": cjk_count >= 3,
        "picture_shape_present": picture_shapes >= 1,
        "embedded_media_present": len(media) >= 1,
        "expected_font_declared": args.expected_font in fonts or args.expected_font in raw_xml,
    }
    result = {
        "schema_version": 1,
        "artifact": args.pptx.name,
        "sha256": sha256(args.pptx),
        "bytes": args.pptx.stat().st_size,
        "slide_count": len(deck.slides),
        "slide_size_emu": [width, height],
        "aspect_ratio": aspect,
        "text_shape_count": text_shapes,
        "editable_cjk_character_count": cjk_count,
        "picture_shape_count": picture_shapes,
        "embedded_media_count": len(media),
        "declared_run_fonts": dict(fonts),
        "checks": checks,
        "pass": all(checks.values()),
    }
    args.json_out.parent.mkdir(parents=True, exist_ok=True)
    args.json_out.write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")
    if not result["pass"]:
        raise SystemExit(1)


if __name__ == "__main__":
    main()
