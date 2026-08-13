from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


FONT = "Microsoft YaHei"
NAVY = RGBColor(7, 29, 45)
TEAL = RGBColor(66, 211, 164)
WHITE = RGBColor(244, 251, 248)
MUTED = RGBColor(159, 233, 211)


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: int, color=WHITE, bold=False):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_background(slide):
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = NAVY


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--image", type=Path, required=True)
    parser.add_argument("--output", type=Path, required=True)
    args = parser.parse_args()

    deck = Presentation()
    deck.slide_width = Inches(13.333333)
    deck.slide_height = Inches(7.5)
    blank = deck.slide_layouts[6]

    slide = deck.slides.add_slide(blank)
    add_background(slide)
    add_text(slide, "本地智能工作台", 0.9, 1.15, 11.5, 0.9, 42, bold=True)
    add_text(slide, "Windows 可行性辅助验证 · 非 OpenCode+Skill 门禁产物", 0.9, 2.2, 11.5, 0.6, 21, MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(3.1), Inches(2.2), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()
    add_text(slide, "16:9 · 中文可编辑 · 本地图片", 0.9, 4.55, 11.5, 0.6, 24)

    slide = deck.slides.add_slide(blank)
    add_background(slide)
    add_text(slide, "私有字体与图片路径", 0.8, 0.55, 11.8, 0.7, 32, bold=True)
    slide.shapes.add_picture(str(args.image), Inches(0.8), Inches(1.55), width=Inches(7.15), height=Inches(4.02))
    add_text(slide, "隔离 Python", 8.45, 1.75, 3.8, 0.5, 25, TEAL, True)
    add_text(slide, "私有目录加载字体\n生成左侧 PNG", 8.45, 2.45, 3.9, 1.35, 19)
    add_text(slide, "右侧文字仍为原生文本\n可在 PowerPoint 中编辑", 8.45, 4.2, 3.9, 1.2, 19, MUTED)

    slide = deck.slides.add_slide(blank)
    add_background(slide)
    add_text(slide, "确定性检查结果", 0.8, 0.6, 11.8, 0.7, 32, bold=True)
    rows = [
        ("01", "OOXML / ZIP", "包结构与 CRC 可独立解析"),
        ("02", "python-pptx", "可编辑中文、图片与 16:9 尺寸"),
        ("03", "PowerPoint COM", "桌面 Office 实际打开并导出预览"),
    ]
    y = 1.65
    for number, title, detail in rows:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y), Inches(11.75), Inches(1.25))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(13, 53, 80)
        card.line.color.rgb = RGBColor(32, 91, 112)
        add_text(slide, number, 1.1, y + 0.27, 0.9, 0.45, 22, TEAL, True)
        add_text(slide, title, 2.05, y + 0.22, 3.15, 0.48, 20, WHITE, True)
        add_text(slide, detail, 5.35, y + 0.24, 6.45, 0.55, 18, MUTED)
        y += 1.55

    args.output.parent.mkdir(parents=True, exist_ok=True)
    deck.save(args.output)


if __name__ == "__main__":
    main()
